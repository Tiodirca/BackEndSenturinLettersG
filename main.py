import socket
from pptx import Presentation
import os
from flask import Flask, send_from_directory
from flask import request
from waitress import serve

# pip install Flask
# pip install python-pptx

gerarArquivo = Flask("name")

nome_letra = ""
letra_completa = []
formato_arquivo = ".ppt"


@gerarArquivo.route("/")
def gerar():
    try:
        return "<p>Tela Inicial.Você está tentando acessar o back-end " \
               "do Senturion Letters G. funções disponiveis somente atraves da aplicação.</p>"
    except Exception as e:
        print(e)
        return f"<p>Erro ao acessar f'{e}'</p>"


@gerarArquivo.route("/chamarBaixarArquivo")
def chamar_baixar_arquivo():
    try:
        return baixar_arquivo(nome_letra)
    except Exception as e:
        print(e)
        return f"<p>erro ao chamar baixar arquivo : '{e}'</p>"


def chamar_criar_arquivo(nomeletra, tipo_modelo):
    # metodo responsavel por chamar a criacao do arquivo
    global letra_completa, formato_arquivo
    if 'modelo_geral' in tipo_modelo:
        # passando o caminho do modelo que sera utilizado para gerar os slides
        caminho = "modelos_slides/modelo_geral.pptx"
        prs = Presentation(caminho)
    else:
        # passando o caminho do modelo que sera utilizado para gerar os slides
        caminho = "modelos_slides/modelo_geracao_fire.pptx"
        prs = Presentation(caminho)
    for verso in letra_completa:
        print(verso + "asdasdszxczxc")
        adicionar_slides(prs, nomeletra, verso)
    prs.save(nomeletra + formato_arquivo)
    # limpando valor da lista
    letra_completa = []


def baixar_arquivo(nome_arquivo):
    # metodo responsavel por baixar o arquivo gerado
    try:
        global formato_arquivo
        caminho_absoluto_arquivo_python = os.path.abspath(__file__)
        # pegando o diretorio baseado no caminho absoluto
        diretorio_src = os.path.dirname(caminho_absoluto_arquivo_python)
        arquivo = nome_arquivo + formato_arquivo
        return send_from_directory(diretorio_src,
                                   arquivo, as_attachment=False)
    except Exception as e:
        print(e)
        return f"<p>erro ao baixar arquivo '{nome_arquivo}' : '{e}'</p>"


def adicionar_slides(prs, titulo, estrofe):
    # metodo responsavel por pegar e adicionar um
    # pagina ao arquivo de slide
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    subtitulo = slide.placeholders[1]
    title.text = titulo
    subtitulo.text = estrofe.replace("pnj", "\n").replace("</p>", "").replace("\r", "")
    print("sa" + subtitulo.text)


@gerarArquivo.route("/pegarValores", methods=['POST'])
def pegar_valores():
    # metodo responsavel por iniciar o processo de criacao do arquivo
    # pegando os valores
    try:
        # a variavel recebe o valor passado atraves do map
        global nome_letra, letra_completa
        tamanholista = request.form.get("tamanhoLista")
        tipo_modelo = request.form.get("modelo_slide")
        nome_letra = request.form.get("nome_letra")
        # fazendo interacao com o tamanho recebido
        for index in range(int(tamanholista)):
            # variavel recebe os valores passados atraves do map
            # recebendo como parametro o index corresponde
            estrofe = request.form.get(f"versos[{index}]")
            # adicionando valores corresponte a uma lista
            letra_completa.append(estrofe)
            print(estrofe + "sadas")
            print(letra_completa)
        chamar_criar_arquivo(nome_letra, tipo_modelo)
        return "<p>sucesso ao pegar valores</p>"
    except Exception as e:
        print(e)
        return f"<p>erro ao pegar valores f'{e}'</p>"


@gerarArquivo.route("/excluirArquivo", methods=['POST'])
def excluir_arquivo_diretorio():
    global formato_arquivo
    # metodo para excluir o arquivo gerado
    arquivo = request.form.get('arquivo')
    nome_arquivo = arquivo + formato_arquivo
    # pegando o caminho absoluto do arquivo
    caminho_absoluto_arquivo_python = os.path.abspath(__file__)
    # pegando o diretorio usando o caminho absoluto
    diretorio_src = os.path.dirname(caminho_absoluto_arquivo_python)
    # listando todos os arquivos
    diretorio = os.listdir(diretorio_src)
    try:
        # verificando se o nome do arquivo corresponde a 
        # algum arquivo contido no diretorio
        for file in diretorio:
            if file == nome_arquivo:
                # removendo o arquivo
                os.remove(file)
        return "<p>sucesso ao excluir</p>"
    except Exception as e:
        print(e)
        return f"<p>erro ao excluir '{e}'</p>"


def obter_ip():
    # metodo para obter ip da maquina para testes
    try:
        # obtendo ip
        ip = socket.gethostbyname(socket.gethostname())
        print(f"Ip(protocolo de rede) da Maquina : {ip}")
        # retornando o arquivo
        return ip
    except Exception as e:
        print(e)
        return f"<p>Erro '{e}'</p>"


if __name__ == '__main__':
    # port = int(os.environ.get("PORT", 5000))
    # gerarArquivo.run(host='0.0.0.0', port=port, debug=True)
    print("Necessario está tela estar aberta para "
          "\ngerar e baixar o arquivo criado no programa"
          "\nSenturion Letters G.")
    serve(gerarArquivo, host=obter_ip(), port=5000)
    # gerarArquivo.run(host=obter_ip(), port=5000, debug=True),
